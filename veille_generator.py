import os
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

def generer_veille(date_str):
    # Simulation de la génération de la veille juridique
    pptx_filename = f"veille_droit_social_{date_str}.pptx"
    script_filename = f"script_presentation_{date_str}.txt"

    prs = Presentation()

    # Création d'une diapositive exemple
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = "Exemple de décision"

    # Sauvegarde du fichier PowerPoint
    prs.save(f"outputs/{pptx_filename}")

    # Génération du script texte
    with open(f"outputs/{script_filename}", "w") as f:
        f.write("Voici un exemple de texte pour la présentation.")

    return pptx_filename, script_filename
